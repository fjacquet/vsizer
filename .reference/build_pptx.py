"""Cluster utilization PPTX — narrative: 'serveurs qui ronronnent → resize en GHz'."""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

XLSX = "/Users/fjacquet/Library/CloudStorage/OneDrive-Home/WIP-sizing/Classeur2.xlsx"
OUT = "/tmp/cluster-pptx/cluster_utilization.pptx"

# Midnight Executive + accents
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x29, 0x5C)
GREY = RGBColor(0x6B, 0x72, 0x91)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFC)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xE0, 0x7B, 0x00)
RED = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xF9, 0xB9, 0x35)
TEAL = RGBColor(0x02, 0x88, 0x9C)


def usage_color(pct):
    if pct < 40:
        return GREEN
    if pct < 70:
        return ORANGE
    return RED


# ---------- Data + GHz analytics ---------------------------------------------
vhost = pd.read_excel(XLSX, sheet_name="vHost")
vinfo = pd.read_excel(XLSX, sheet_name="vInfo")

# Per-host GHz available and actually consumed
vhost["ghz_phys"] = vhost["Speed"] * vhost["# Cores"] / 1000.0  # MHz × cores → GHz
vhost["ghz_used"] = vhost["ghz_phys"] * vhost["CPU usage %"] / 100.0
# A common sizing assumption: 1 vCPU ≈ 1 core's worth of GHz (host clock)
vhost["ghz_alloc_assumed"] = vhost["# vCPUs"] * vhost["Speed"] / 1000.0

agg = vhost.groupby("Cluster").agg(
    nb_hosts=("Host", "count"),
    speed_mhz=("Speed", "mean"),
    cpu_avg=("CPU usage %", "mean"),
    mem_avg=("Memory usage %", "mean"),
    cpu_max=("CPU usage %", "max"),
    mem_max=("Memory usage %", "max"),
    cpu_min=("CPU usage %", "min"),
    mem_min=("Memory usage %", "min"),
    total_cores=("# Cores", "sum"),
    total_mem_mb=("# Memory", "sum"),
    total_vcpus=("# vCPUs", "sum"),
    total_vram_mb=("vRAM", "sum"),
    ghz_phys=("ghz_phys", "sum"),
    ghz_used=("ghz_used", "sum"),
    ghz_alloc_assumed=("ghz_alloc_assumed", "sum"),
).reset_index()

vinfo_on = vinfo[vinfo["Powerstate"].astype(str).str.contains("on", case=False, na=False)]
vi = vinfo_on.groupby("Cluster").agg(
    vms_on=("VM", "count"),
    vcpus_alloc=("CPUs", "sum"),
    vram_alloc_mib=("Memory", "sum"),
    active_mem_mib=("Active Memory", "sum"),
).reset_index()
data = agg.merge(vi, on="Cluster", how="left").fillna(0)
data["mhz_per_vcpu_real"] = data.apply(
    lambda r: (r["ghz_used"] * 1000 / r["total_vcpus"]) if r["total_vcpus"] else 0, axis=1
)
data["headroom_ghz"] = data["ghz_phys"] - data["ghz_used"]
data["headroom_pct"] = data["headroom_ghz"] / data["ghz_phys"] * 100
data = data.sort_values("Cluster").reset_index(drop=True)

# Global aggregates
TOT_GHZ_PHYS = data["ghz_phys"].sum()
TOT_GHZ_USED = data["ghz_used"].sum()
TOT_VCPU = data["total_vcpus"].sum()
TOT_VMS = data["vms_on"].sum()
TOT_HOSTS = data["nb_hosts"].sum()
GLOBAL_CPU_AVG = TOT_GHZ_USED / TOT_GHZ_PHYS * 100
GLOBAL_MHZ_PER_VCPU = TOT_GHZ_USED * 1000 / TOT_VCPU if TOT_VCPU else 0
GLOBAL_MEM_AVG = (vhost["Memory usage %"] * vhost["# Memory"]).sum() / vhost["# Memory"].sum()

# ---------- Presentation -----------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK_TEXT, font="Calibri", align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if rounded:
        # Reduce corner radius
        sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def fmt_mb(mb):
    mb = float(mb)
    if mb >= 1024 * 1024:
        return f"{mb / 1024 / 1024:.1f} TB"
    if mb >= 1024:
        return f"{mb / 1024:.1f} GB"
    return f"{mb:.0f} MB"


def fmt_int(x):
    return f"{int(x):,}".replace(",", " ")


def fmt_ghz(ghz):
    return f"{ghz:,.0f} GHz".replace(",", " ")


# ---------- Slide 1: Hero / message clé --------------------------------------
def slide_hero():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    # Decorative gold strip
    add_rect(s, 0, Inches(6.75), SW, Inches(0.18), GOLD)

    add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.45),
             "ANALYSE DE CAPACITÉ — VMware",
             size=14, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(2.0),
             "Vos serveurs\nronronnent.",
             size=88, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.6),
             "Sizing en vCPU = surdimensionnement.  Sizing en GHz = vérité.",
             size=24, italic=True, color=ICE)

    # Three big numbers
    nums = [
        (f"{GLOBAL_CPU_AVG:.0f} %", "CPU moyen utilisé\nsur l'ensemble des clusters"),
        (f"{TOT_GHZ_PHYS - TOT_GHZ_USED:,.0f}".replace(",", " ") + " GHz",
         "de capacité CPU\ninutilisée"),
        (f"{GLOBAL_MHZ_PER_VCPU:.0f} MHz", "consommés en moyenne\npar vCPU alloué"),
    ]
    nx = Inches(0.7)
    nw = Inches(4.0)
    for big, small in nums:
        add_rect(s, nx, Inches(5.0), nw, Inches(1.55),
                 RGBColor(0x2A, 0x35, 0x80), rounded=True)
        add_text(s, nx + Inches(0.2), Inches(5.1), nw - Inches(0.4), Inches(0.7),
                 big, size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(s, nx + Inches(0.2), Inches(5.85), nw - Inches(0.4), Inches(0.7),
                 small, size=12, color=ICE, align=PP_ALIGN.CENTER)
        nx += nw + Inches(0.15)


# ---------- Slide 2: Le message — vCPU vs GHz --------------------------------
def slide_message():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, WHITE)
    add_rect(s, 0, 0, SW, Inches(1.2), NAVY)
    add_text(s, Inches(0.5), Inches(0.25), Inches(12), Inches(0.55),
             "Le constat : sizing en vCPU ≠ consommation réelle",
             size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(0.85), Inches(12), Inches(0.3),
             f"{fmt_int(TOT_VCPU)} vCPU alloués sur {fmt_int(int(data['total_cores'].sum()))} "
             f"cores physiques — pour {GLOBAL_CPU_AVG:.0f} % d'utilisation moyenne",
             size=13, color=ICE)

    # Two side-by-side panels
    left_x = Inches(0.6)
    right_x = Inches(6.93)
    pw = Inches(5.8)
    py = Inches(1.6)
    ph = Inches(5.3)

    # LEFT — vCPU sizing (the problem)
    add_rect(s, left_x, py, pw, ph, LIGHT_BG, rounded=True)
    add_rect(s, left_x, py, Inches(0.18), ph, RED)  # red rail
    add_text(s, left_x + Inches(0.4), py + Inches(0.25), pw - Inches(0.6), Inches(0.4),
             "❌  SIZING vCPU", size=11, bold=True, color=RED)
    add_text(s, left_x + Inches(0.4), py + Inches(0.65), pw - Inches(0.6), Inches(0.7),
             "1 vCPU =\n1 cœur entier ?", size=30, bold=True, color=NAVY)
    add_text(s, left_x + Inches(0.4), py + Inches(2.1), pw - Inches(0.6), Inches(1.4),
             f"{fmt_int(TOT_VCPU)}",
             size=72, bold=True, color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left_x + Inches(0.4), py + Inches(3.55), pw - Inches(0.6), Inches(0.35),
             "vCPU alloués", size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, left_x + Inches(0.4), py + Inches(4.15), pw - Inches(0.6), Inches(0.9),
             "Capacité réservée par défaut : "
             f"{TOT_VCPU * data['speed_mhz'].mean() / 1000:,.0f} GHz".replace(",", " ")
             + "\n→ provisioning massif, latent et inutilisé.",
             size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # RIGHT — GHz sizing (the truth)
    add_rect(s, right_x, py, pw, ph, LIGHT_BG, rounded=True)
    add_rect(s, right_x, py, Inches(0.18), ph, GREEN)
    add_text(s, right_x + Inches(0.4), py + Inches(0.25), pw - Inches(0.6), Inches(0.4),
             "✅  SIZING GHz", size=11, bold=True, color=GREEN)
    add_text(s, right_x + Inches(0.4), py + Inches(0.65), pw - Inches(0.6), Inches(0.7),
             "Ce que les VMs\nconsomment vraiment", size=30, bold=True, color=NAVY)
    add_text(s, right_x + Inches(0.4), py + Inches(2.1), pw - Inches(0.6), Inches(1.4),
             f"{TOT_GHZ_USED:,.0f}".replace(",", " ") + " GHz",
             size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, right_x + Inches(0.4), py + Inches(3.55), pw - Inches(0.6), Inches(0.35),
             f"sur {fmt_ghz(TOT_GHZ_PHYS)} disponibles",
             size=14, color=GREY, align=PP_ALIGN.CENTER)
    add_text(s, right_x + Inches(0.4), py + Inches(4.15), pw - Inches(0.6), Inches(0.9),
             f"Soit {GLOBAL_MHZ_PER_VCPU:.0f} MHz par vCPU en moyenne — "
             f"~ {GLOBAL_MHZ_PER_VCPU / data['speed_mhz'].mean() * 100:.0f} % d'un cœur.\n"
             "→ resize possible. Cibler la capacité réelle.",
             size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.5), Inches(7.05), Inches(12), Inches(0.3),
             "Source : Classeur2.xlsx — vHost.Speed × #Cores × CPU usage %  ·  vInfo.CPUs (powered on)",
             size=9, color=GREY)


# ---------- Slide 3: Overview chart ------------------------------------------
def slide_overview():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, WHITE)
    add_rect(s, 0, 0, SW, Inches(1.0), NAVY)
    add_text(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             "Vue d'ensemble — utilisation CPU & RAM par cluster",
             size=26, bold=True, color=WHITE)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.3),
             "Moyenne pondérée des hôtes  ·  marqueur or = pic max  ·  taille = capacité physique",
             size=11, color=ICE)

    # Layout grid: 18 clusters
    n = len(data)
    top = Inches(1.25)
    bottom_pad = Inches(0.5)
    row_h = (SH - top - bottom_pad) / n

    name_x = Inches(0.45)
    name_w = Inches(1.35)
    info_x = Inches(1.85)
    info_w = Inches(1.35)
    bar_zone_x = Inches(3.30)
    bar_zone_w = Inches(6.6)
    pct_x = bar_zone_x + bar_zone_w + Inches(0.15)
    pct_w = Inches(1.5)
    head_x = pct_x + pct_w + Inches(0.05)
    head_w = Inches(1.55)

    bar_h_pt = 9
    bar_h = Emu(int(Pt(bar_h_pt).emu))

    # Column headers
    htop = top - Inches(0.22)
    add_text(s, name_x, htop, name_w, Inches(0.2), "Cluster",
             size=8, bold=True, color=GOLD)
    add_text(s, info_x, htop, info_w, Inches(0.2), "Hôtes / VMs",
             size=8, bold=True, color=GOLD)
    add_text(s, bar_zone_x, htop, bar_zone_w, Inches(0.2),
             "0%   ·   utilisation hôtes (haut: CPU, bas: RAM)   ·   100%",
             size=8, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, pct_x, htop, pct_w, Inches(0.2), "Moyenne / Pic",
             size=8, bold=True, color=GOLD)
    add_text(s, head_x, htop, head_w, Inches(0.2), "Marge libérable",
             size=8, bold=True, color=GOLD)

    for i, r in data.iterrows():
        y = top + row_h * i
        if i % 2 == 0:
            add_rect(s, Inches(0.3), y, SW - Inches(0.6),
                     row_h - Emu(int(Pt(1).emu)), LIGHT_BG)
        # Cluster name
        add_text(s, name_x, y, name_w, row_h, r["Cluster"],
                 size=11, bold=True, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        # Sub-info
        add_text(s, info_x, y, info_w, row_h,
                 f"{int(r['nb_hosts'])} hôtes\n{int(r['vms_on'])} VMs",
                 size=8, color=GREY, anchor=MSO_ANCHOR.MIDDLE)

        # Compute centered bar Y positions: 2 bars + small gap
        gap = Emu(int(Pt(2).emu))
        total_bars = bar_h * 2 + gap
        cpu_y = y + (row_h - total_bars) / 2
        mem_y = cpu_y + bar_h + gap

        # Tiny labels left of bars
        label_w = Inches(0.25)
        add_text(s, bar_zone_x - label_w - Emu(int(Pt(2).emu)),
                 cpu_y - Emu(int(Pt(2).emu)),
                 label_w, bar_h + Emu(int(Pt(4).emu)),
                 "CPU", size=7, bold=True, color=GREY,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(s, bar_zone_x - label_w - Emu(int(Pt(2).emu)),
                 mem_y - Emu(int(Pt(2).emu)),
                 label_w, bar_h + Emu(int(Pt(4).emu)),
                 "RAM", size=7, bold=True, color=GREY,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

        # CPU bar
        add_rect(s, bar_zone_x, cpu_y, bar_zone_w, bar_h, ICE)
        cw = int(bar_zone_w * max(0, min(r["cpu_avg"], 100)) / 100)
        if cw > 0:
            add_rect(s, bar_zone_x, cpu_y, cw, bar_h, usage_color(r["cpu_avg"]))
        if r["cpu_max"] > 0:
            mxp = bar_zone_x + int(bar_zone_w * r["cpu_max"] / 100) - Emu(int(Pt(1).emu))
            add_rect(s, mxp, cpu_y - Emu(int(Pt(1.5).emu)),
                     Emu(int(Pt(2).emu)), bar_h + Emu(int(Pt(3).emu)), GOLD)

        # MEM bar
        add_rect(s, bar_zone_x, mem_y, bar_zone_w, bar_h, ICE)
        mw = int(bar_zone_w * max(0, min(r["mem_avg"], 100)) / 100)
        if mw > 0:
            add_rect(s, bar_zone_x, mem_y, mw, bar_h, usage_color(r["mem_avg"]))
        if r["mem_max"] > 0:
            mxp = bar_zone_x + int(bar_zone_w * r["mem_max"] / 100) - Emu(int(Pt(1).emu))
            add_rect(s, mxp, mem_y - Emu(int(Pt(1.5).emu)),
                     Emu(int(Pt(2).emu)), bar_h + Emu(int(Pt(3).emu)), GOLD)

        # Percentages on the right
        add_text(s, pct_x, y, pct_w, row_h / 2,
                 f"CPU  {r['cpu_avg']:.0f}%   pic {r['cpu_max']:.0f}%",
                 size=9, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, pct_x, y + row_h / 2, pct_w, row_h / 2,
                 f"RAM  {r['mem_avg']:.0f}%   pic {r['mem_max']:.0f}%",
                 size=9, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        # Headroom
        add_text(s, head_x, y, head_w, row_h,
                 f"+{r['headroom_ghz']:.0f} GHz  ({r['headroom_pct']:.0f}%)",
                 size=10, bold=True, italic=True, color=GREEN,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Legend footer
    ly = SH - Inches(0.35)
    legend = [("< 40%", GREEN), ("40-70%", ORANGE), ("≥ 70%", RED), ("Pic max", GOLD)]
    add_text(s, Inches(0.45), ly, Inches(0.8), Inches(0.3),
             "Légende :", size=9, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    lx = Inches(1.25)
    for txt, col in legend:
        add_rect(s, lx, ly + Inches(0.06), Inches(0.18), Inches(0.18), col)
        add_text(s, lx + Inches(0.22), ly, Inches(1.6), Inches(0.3), txt,
                 size=9, color=DARK_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        lx += Inches(1.4)


# ---------- Per-cluster slide ------------------------------------------------
def slide_cluster(r):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, WHITE)
    # Left rail
    add_rect(s, 0, 0, Inches(0.35), SH, NAVY)
    # Header
    add_rect(s, Inches(0.35), 0, SW - Inches(0.35), Inches(1.15), NAVY)
    add_text(s, Inches(0.7), Inches(0.18), Inches(10), Inches(0.65),
             r["Cluster"], size=34, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(0.72), Inches(11), Inches(0.35),
             f"{int(r['nb_hosts'])} hôtes  ·  {int(r['vms_on'])} VMs allumées  ·  "
             f"{int(r['total_cores'])} cores phys. ({r['speed_mhz'] / 1000:.2f} GHz/core)  ·  "
             f"{fmt_mb(r['total_mem_mb'])} RAM",
             size=12, color=ICE)

    # ---- ROW 1: KPI cards (4 small) ----
    card_y = Inches(1.35)
    card_h = Inches(1.05)
    cards = [
        (f"{r['cpu_avg']:.0f}%", "CPU moyen", usage_color(r["cpu_avg"])),
        (f"{r['mem_avg']:.0f}%", "RAM moyenne", usage_color(r["mem_avg"])),
        (f"{r['ghz_used']:,.0f}".replace(",", " ") + " / " +
         f"{r['ghz_phys']:,.0f}".replace(",", " "), "GHz utilisés / phys.", NAVY),
        (f"{r['mhz_per_vcpu_real']:.0f} MHz", "réels par vCPU alloué", TEAL),
    ]
    cx = Inches(0.7)
    cw = Inches(2.95)
    gap = Inches(0.15)
    for big, small, col in cards:
        add_rect(s, cx, card_y, cw, card_h, LIGHT_BG, rounded=True)
        add_rect(s, cx, card_y, Inches(0.12), card_h, col)
        add_text(s, cx + Inches(0.25), card_y + Inches(0.1),
                 cw - Inches(0.3), Inches(0.6), big,
                 size=26, bold=True, color=col)
        add_text(s, cx + Inches(0.25), card_y + Inches(0.66),
                 cw - Inches(0.3), Inches(0.35), small,
                 size=11, color=GREY)
        cx += cw + gap

    # ---- ROW 2 : CPU & RAM utilization blocks ----
    block_y = Inches(2.6)
    block_h = Inches(2.1)
    block_w = Inches(6.05)

    def draw_util_block(x, label, avg, mx, mn, sub_top, sub_bot):
        add_rect(s, x, block_y, block_w, block_h, LIGHT_BG, rounded=True)
        add_text(s, x + Inches(0.3), block_y + Inches(0.18),
                 block_w - Inches(0.6), Inches(0.35), label,
                 size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), block_y + Inches(0.5),
                 block_w - Inches(0.6), Inches(0.3), sub_top,
                 size=10, color=GREY)
        # Big bar
        bar_y = block_y + Inches(0.95)
        bar_h = Inches(0.32)
        bar_x = x + Inches(0.3)
        bar_w = block_w - Inches(0.6)
        add_rect(s, bar_x, bar_y, bar_w, bar_h, ICE, rounded=True)
        fill_w = int(bar_w * max(0, min(avg, 100)) / 100)
        if fill_w > 0:
            add_rect(s, bar_x, bar_y, fill_w, bar_h, usage_color(avg), rounded=True)
        # max marker line
        if mx > 0:
            mxp = bar_x + int(bar_w * mx / 100) - Emu(int(Pt(1.5).emu))
            add_rect(s, mxp, bar_y - Inches(0.08),
                     Emu(int(Pt(3).emu)), bar_h + Inches(0.16), GOLD)
        # 0/100 labels
        add_text(s, bar_x, bar_y + bar_h + Inches(0.05), bar_w, Inches(0.2),
                 "0%", size=8, color=GREY)
        add_text(s, bar_x, bar_y + bar_h + Inches(0.05), bar_w, Inches(0.2),
                 "100%", size=8, color=GREY, align=PP_ALIGN.RIGHT)
        # Min/avg/max strip
        sy = block_y + Inches(1.55)
        sw = bar_w / 3
        for i, (lab, val, c) in enumerate([
            ("Min", f"{mn:.0f}%", GREY),
            ("Moy", f"{avg:.1f}%", usage_color(avg)),
            ("Max", f"{mx:.0f}%", GREY),
        ]):
            sx = bar_x + sw * i
            add_text(s, sx, sy, sw, Inches(0.25), lab,
                     size=9, color=GREY, align=PP_ALIGN.CENTER)
            add_text(s, sx, sy + Inches(0.2), sw, Inches(0.3), val,
                     size=14, bold=True, color=c, align=PP_ALIGN.CENTER)

    cpu_top = (
        f"{r['ghz_used']:,.0f}".replace(",", " ")
        + f" GHz consommés sur {r['ghz_phys']:,.0f}".replace(",", " ") + " GHz"
    )
    mem_top = (
        f"{fmt_mb(r['total_mem_mb'] * r['mem_avg'] / 100)} consommés "
        f"sur {fmt_mb(r['total_mem_mb'])}"
    )
    draw_util_block(Inches(0.7), "CPU — utilisation moyenne",
                    r["cpu_avg"], r["cpu_max"], r["cpu_min"], cpu_top, "")
    draw_util_block(Inches(0.7) + block_w + Inches(0.4), "RAM — utilisation moyenne",
                    r["mem_avg"], r["mem_max"], r["mem_min"], mem_top, "")

    # ---- ROW 3 : The "purring" panel — vCPU vs GHz used ----
    p_y = Inches(4.85)
    p_h = Inches(1.95)
    # Big panel full width
    add_rect(s, Inches(0.7), p_y, SW - Inches(1.4), p_h, NAVY, rounded=True)
    add_text(s, Inches(0.95), p_y + Inches(0.15), Inches(8), Inches(0.4),
             "💡  RESIZE EN GHZ — POTENTIEL", size=12, bold=True, color=GOLD)
    add_text(s, Inches(0.95), p_y + Inches(0.5), Inches(11.6), Inches(0.45),
             f"Ce cluster ronronne à {r['cpu_avg']:.0f} % de CPU.",
             size=22, bold=True, color=WHITE)

    # Inline comparison: vCPU·MHz_assumed vs GHz_used
    metrics = [
        ("vCPU alloués", fmt_int(r["total_vcpus"])),
        ("Capacité réservée\n(1 vCPU ≈ 1 core)",
         f"{r['ghz_alloc_assumed']:,.0f}".replace(",", " ") + " GHz"),
        ("Réellement consommé",
         f"{r['ghz_used']:,.0f}".replace(",", " ") + " GHz"),
        ("Marge libérable",
         f"{r['headroom_ghz']:,.0f}".replace(",", " ") + f" GHz ({r['headroom_pct']:.0f} %)"),
    ]
    mx = Inches(0.95)
    mw = (SW - Inches(2.0)) / len(metrics)
    for lab, val in metrics:
        add_text(s, mx, p_y + Inches(1.15), mw, Inches(0.32),
                 lab, size=9, color=ICE)
        add_text(s, mx, p_y + Inches(1.45), mw, Inches(0.45),
                 val, size=18, bold=True, color=GOLD)
        mx += mw

    # Footer
    add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
             "Source : Classeur2.xlsx — vHost (CPU/RAM usage %, Speed × Cores) + vInfo (vCPUs, Memory)",
             size=9, color=GREY)


# ---------- Closing slide ----------------------------------------------------
def slide_recommend():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(6.75), SW, Inches(0.18), GOLD)

    add_text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.4),
             "RECOMMANDATIONS", size=14, bold=True, color=GOLD)
    add_text(s, Inches(0.7), Inches(1.15), Inches(12), Inches(1.2),
             "Du sizing vCPU au sizing GHz", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
             f"Marge disponible globale : {fmt_ghz(TOT_GHZ_PHYS - TOT_GHZ_USED)} sur {fmt_ghz(TOT_GHZ_PHYS)}",
             size=18, italic=True, color=ICE)

    actions = [
        ("1", "Mesurer en MHz consommés",
         f"Sur l'estate, chaque vCPU consomme ~ {GLOBAL_MHZ_PER_VCPU:.0f} MHz en moyenne — pas un cœur entier."),
        ("2", "Resize les VMs sur la conso réelle",
         "Cibler le 95ᵉ percentile MHz observé, pas la déclaration vCPU initiale."),
        ("3", "Consolider les clusters sous-utilisés",
         "Plusieurs clusters tournent sous 15 % CPU : candidats à la fusion ou à la décommission."),
        ("4", "Refresh : sizer en GHz, pas en cores",
         "Le sizer Dell AI/VMware accepte une cible GHz — fini la course aux cores 'au cas où'."),
    ]
    ay = Inches(3.2)
    for num, title, desc in actions:
        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), ay,
                                    Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOLD
        circle.line.fill.background()
        circle.shadow.inherit = False
        tf = circle.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = NAVY

        add_text(s, Inches(1.45), ay - Inches(0.02), Inches(11), Inches(0.35),
                 title, size=18, bold=True, color=WHITE)
        add_text(s, Inches(1.45), ay + Inches(0.32), Inches(11), Inches(0.4),
                 desc, size=12, color=ICE)
        ay += Inches(0.85)

    add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
             "Source : Classeur2.xlsx — vInfo + vHost",
             size=10, color=ICE)


# ---------- Build ------------------------------------------------------------
slide_hero()
slide_message()
slide_overview()
for _, row in data.iterrows():
    slide_cluster(row)
slide_recommend()

prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
